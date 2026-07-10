# -*- coding: utf-8 -*-
"""
하버스카페 — 손님 리뷰(VoC) 엑셀 리포트 + 경쟁사 분석 PPT 생성기.
- 엑셀: cafe_reviews.csv(원천 데이터)를 읽어, 실제 Excel 수식(COUNTIF/AVERAGEIF/COUNTIFS)과
  차트로 분석. 값 하드코딩(박제) 아님 — 원본 시트를 참조하는 수식으로 계산.
- PPT: competitors.md 내용을 표지/개요/비교표/팔로워차트/차별화/추천액션 슬라이드로.
"""
import csv, os
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from openpyxl.utils import get_column_letter

HERE = os.path.dirname(os.path.abspath(__file__))
CSV = os.path.join(HERE, "cafe_reviews.csv")

# ---------- 공통 스타일 ----------
NAVY = "1F3A5F"; AMBER = "B45309"; LIGHT = "F5EFE6"; RED = "C0392B"
thin = Side(style="thin", color="D9D2C5")
BORDER = Border(left=thin, right=thin, top=thin, bottom=thin)
HFONT = Font(name="맑은 고딕", bold=True, color="FFFFFF", size=11)
BFONT = Font(name="맑은 고딕", size=10)
TITLE = Font(name="맑은 고딕", bold=True, size=14, color=NAVY)

def style_header(ws, row, ncols, fill=NAVY):
    for c in range(1, ncols + 1):
        cell = ws.cell(row=row, column=c)
        cell.font = HFONT; cell.fill = PatternFill("solid", fgColor=fill)
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = BORDER

# ================= 1) 원천 데이터 로드 =================
rows = []
with open(CSV, encoding="utf-8") as f:
    for r in csv.DictReader(f):
        rows.append(r)
n = len(rows)                      # 리뷰 수
last = n + 1                       # 원본 데이터 마지막 행(헤더 포함)

# ================= 2) 엑셀 리포트 =================
wb = Workbook()

# --- 시트 A: 리뷰 원본 ---
raw = wb.active; raw.title = "리뷰원본"
heads = ["날짜", "플랫폼", "별점", "리뷰", "테마"]
raw.append(heads); style_header(raw, 1, len(heads))
for r in rows:
    raw.append([r["날짜"], r["플랫폼"], int(r["별점"]), r["리뷰"], r["테마"]])
for i in range(2, last + 1):
    for c in range(1, 6):
        raw.cell(row=i, column=c).border = BORDER
        raw.cell(row=i, column=c).font = BFONT
    raw.cell(row=i, column=3).alignment = Alignment(horizontal="center")
widths = [12, 9, 7, 40, 10]
for i, w in enumerate(widths, 1):
    raw.column_dimensions[get_column_letter(i)].width = w
raw.freeze_panes = "A2"

RC = f"리뷰원본!$C$2:$C${last}"   # 별점 범위
RE = f"리뷰원본!$E$2:$E${last}"   # 테마 범위

# --- 시트 B: 분석 (수식 기반) ---
an = wb.create_sheet("VoC분석")
an["A1"] = "하버스카페 — 손님의 소리(VoC) 분석 리포트"; an["A1"].font = TITLE
an["A2"] = f"원천: cafe_reviews.csv · 리뷰 {n}건 · 네이버/카카오맵/인스타"; an["A2"].font = Font(name="맑은 고딕", size=9, color="7A7266")

# KPI (실제 수식)
an["A4"] = "핵심 지표"; an["A4"].font = Font(name="맑은 고딕", bold=True, size=12, color=AMBER)
kpis = [
    ("총 리뷰 수", f"=COUNTA(리뷰원본!A2:A{last})"),
    ("평균 별점", f"=ROUND(AVERAGE({RC}),2)"),
    ("긍정률(★4↑)", f"=ROUND(COUNTIF({RC},\">=4\")/COUNTA({RC}),3)"),
    ("부정 리뷰(★2↓)", f"=COUNTIF({RC},\"<=2\")"),
]
for i, (k, v) in enumerate(kpis):
    col = 1 + i * 2
    an.cell(row=5, column=col, value=k).font = HFONT
    an.cell(row=5, column=col).fill = PatternFill("solid", fgColor=AMBER)
    an.cell(row=5, column=col).alignment = Alignment(horizontal="center")
    an.cell(row=5, column=col).border = BORDER
    cell = an.cell(row=6, column=col, value=v)
    cell.font = Font(name="맑은 고딕", bold=True, size=13, color=NAVY)
    cell.alignment = Alignment(horizontal="center"); cell.border = BORDER
    if k == "긍정률(★4↑)": cell.number_format = "0.0%"

# 별점 분포 (COUNTIF)
an["A9"] = "별점 분포"; an["A9"].font = Font(name="맑은 고딕", bold=True, size=12, color=AMBER)
an.append([]) if False else None
an["A10"] = "별점"; an["B10"] = "리뷰 수"; style_header(an, 10, 2)
for s in range(5, 0, -1):
    row = 16 - s  # 별점5 -> 11 ... 별점1 -> 15
    an.cell(row=row, column=1, value=f"★{s}").alignment = Alignment(horizontal="center")
    an.cell(row=row, column=1).border = BORDER; an.cell(row=row, column=1).font = BFONT
    fc = an.cell(row=row, column=2, value=f"=COUNTIF({RC},{s})")
    fc.border = BORDER; fc.font = BFONT; fc.alignment = Alignment(horizontal="center")

chart1 = BarChart(); chart1.type = "col"; chart1.title = "별점 분포"; chart1.height = 6.5; chart1.width = 11
data = Reference(an, min_col=2, min_row=10, max_row=15)
cats = Reference(an, min_col=1, min_row=11, max_row=15)
chart1.add_data(data, titles_from_data=True); chart1.set_categories(cats)
chart1.legend = None
an.add_chart(chart1, "D9")

# 테마별 집계 (COUNTIF + AVERAGEIF) — 피벗 성격
an["A18"] = "테마별 집계 (리뷰 수 · 평균 별점)"; an["A18"].font = Font(name="맑은 고딕", bold=True, size=12, color=AMBER)
an["A19"] = "테마"; an["B19"] = "리뷰 수"; an["C19"] = "평균 별점"; style_header(an, 19, 3)
themes = []
for r in rows:
    if r["테마"] not in themes: themes.append(r["테마"])
# 리뷰 수 많은 순으로 정렬하기 위해 사전 카운트(표시는 수식으로)
themes.sort(key=lambda t: -sum(1 for r in rows if r["테마"] == t))
tstart = 20
for i, t in enumerate(themes):
    row = tstart + i
    an.cell(row=row, column=1, value=t).border = BORDER; an.cell(row=row, column=1).font = BFONT
    c2 = an.cell(row=row, column=2, value=f'=COUNTIF({RE},"{t}")'); c2.border = BORDER; c2.font = BFONT; c2.alignment = Alignment(horizontal="center")
    c3 = an.cell(row=row, column=3, value=f'=ROUND(AVERAGEIF({RE},"{t}",{RC}),1)'); c3.border = BORDER; c3.font = BFONT; c3.alignment = Alignment(horizontal="center")
tend = tstart + len(themes) - 1

chart2 = BarChart(); chart2.type = "bar"; chart2.title = "테마별 리뷰 수"; chart2.height = 7; chart2.width = 11
d2 = Reference(an, min_col=2, min_row=19, max_row=tend)
c2r = Reference(an, min_col=1, min_row=20, max_row=tend)
chart2.add_data(d2, titles_from_data=True); chart2.set_categories(c2r); chart2.legend = None
an.add_chart(chart2, "E18")

# 부정 리뷰(★2↓) 테마 Top3 (COUNTIFS)
nrow = tend + 3
an.cell(row=nrow, column=1, value="부정 리뷰(★2 이하) 테마 Top3").font = Font(name="맑은 고딕", bold=True, size=12, color=RED)
an.cell(row=nrow+1, column=1, value="테마"); an.cell(row=nrow+1, column=2, value="부정 건수")
style_header(an, nrow+1, 2, fill=RED)
neg_sorted = sorted(themes, key=lambda t: -sum(1 for r in rows if r["테마"] == t and int(r["별점"]) <= 2))
for i, t in enumerate(neg_sorted[:3]):
    row = nrow + 2 + i
    an.cell(row=row, column=1, value=t).border = BORDER; an.cell(row=row, column=1).font = BFONT
    c = an.cell(row=row, column=2, value=f'=COUNTIFS({RE},"{t}",{RC},"<=2")'); c.border = BORDER; c.font = BFONT; c.alignment = Alignment(horizontal="center")

# 인사이트 한 줄
insight_row = nrow + 6
an.cell(row=insight_row, column=1, value="📌 인사이트: 긍정은 '맛·작업환경', 부정 1위는 '대기시간(피크 웨이팅)' → 최우선 개선 = 피크타임 대기 단축").font = Font(name="맑은 고딕", bold=True, size=11, color=NAVY)

for col, w in zip("ABCDE", [16, 10, 10, 12, 12]):
    an.column_dimensions[col].width = w

xlsx_path = os.path.join(HERE, "하버스카페_VoC분석.xlsx")
wb.save(xlsx_path)
print("엑셀 저장:", xlsx_path)

# ================= 3) 경쟁사 분석 PPT =================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

NAVYC = RGBColor(0x1F, 0x3A, 0x5F); AMBERC = RGBColor(0xB4, 0x53, 0x09)
GRAY = RGBColor(0x7A, 0x72, 0x66); WHITE = RGBColor(0xFF, 0xFF, 0xFF); REDC = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def add_bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def txt(slide, l, t, w, h, text, size=18, bold=False, color=NAVYC, align=PP_ALIGN.LEFT, font="맑은 고딕"):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color; run.font.name = font
    return tb

def bullets(slide, l, t, w, h, items, size=16):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, (txt_, sub) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = "•  " + txt_
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = NAVYC; r.font.name = "맑은 고딕"
        if sub:
            p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = "     " + sub
            r2.font.size = Pt(size-3); r2.font.color.rgb = GRAY; r2.font.name = "맑은 고딕"
        p.space_after = Pt(6)
    return tb

def title_bar(slide, text):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVYC; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "맑은 고딕"

# 경쟁사 데이터
comp = [
    ("성수로스터스", 5000, "싱글오리진 핸드드립", 9400, "원두 퀄리티", "가격↑·좌석↓"),
    ("데일리빈", 4000, "대용량 라떼", 3200, "가성비·회전", "디저트 약함"),
    ("스윗아워", 4800, "수제 디저트", 13000, "디저트·인스타", "커피 평범"),
    ("하버스카페(우리)", 4500, "흑임자 라떼", 2400, "흑임자·작업·친절", "인지도↓·웨이팅"),
]

# --- 슬라이드 1: 표지 ---
s = prs.slides.add_slide(BLANK); add_bg(s, NAVYC)
band = s.shapes.add_shape(1, 0, Inches(2.7), SW, Inches(0.08)); band.fill.solid(); band.fill.fore_color.rgb = AMBERC; band.line.fill.background()
txt(s, Inches(1), Inches(2.9), Inches(11), Inches(1.2), "하버스카페 경쟁사 분석 보고서", 40, True, WHITE)
txt(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8), "성수동 스페셜티 카페 상권 · 시장조사 → 의사결정", 20, False, RGBColor(0xC9,0xD6,0xE5))
txt(s, Inches(1), Inches(6.4), Inches(11), Inches(0.5), "작성: 하버스카페 · 원천: competitors.md", 13, False, RGBColor(0x9A,0xA9,0xBD))

# --- 슬라이드 2: 시장 개요 ---
s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); title_bar(s, "시장 개요")
bullets(s, Inches(0.7), Inches(1.6), Inches(12), Inches(5), [
    ("반경 500m 내 스페셜티 카페 4곳 · 직장인·디자이너·프리랜서 상권", "평균 아메리카노 4,700원 · 주말 나들이객 + 평일 작업 손님 공존"),
    ("핵심 경쟁축: 시그니처 정체성 · 디저트 · 인스타 인지도 · 작업 환경", "좌석·콘센트 기반 작업 환경은 우리의 명확한 우위"),
    ("상권 공통 약점 = 피크타임 웨이팅", "우리 VoC 리뷰 불만 1위(★2, 대기시간)와 정확히 겹침"),
], size=18)

# --- 슬라이드 3: 경쟁사 비교표 ---
s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); title_bar(s, "경쟁사 비교")
cols = ["카페", "아메리카노", "시그니처", "인스타 팔로워", "강점", "약점"]
trows = len(comp) + 1
tbl_shape = s.shapes.add_table(trows, len(cols), Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.6))
table = tbl_shape.table
for j, h in enumerate(cols):
    cell = table.cell(0, j); cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVYC
    pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    pr.runs[0].font.size = Pt(14); pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = WHITE; pr.runs[0].font.name = "맑은 고딕"
for i, (name, am, sig, fol, st, wk) in enumerate(comp, 1):
    ours = "우리" in name
    vals = [name, f"{am:,}원", sig, f"{fol:,}", st, wk]
    for j, v in enumerate(vals):
        cell = table.cell(i, j); cell.text = v
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]; r.font.size = Pt(12); r.font.name = "맑은 고딕"
        r.font.bold = ours; r.font.color.rgb = AMBERC if ours else NAVYC
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFB,0xF3,0xE6) if ours else WHITE

# --- 슬라이드 4: 인스타 팔로워 격차 (차트) ---
s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); title_bar(s, "인스타 인지도 격차")
cd = CategoryChartData(); cd.categories = [c[0].replace("(우리)", "") for c in comp]
cd.add_series("인스타 팔로워", [c[3] for c in comp])
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.5), Inches(8.2), Inches(5.3), cd)
gframe.chart.has_legend = False
txt(s, Inches(9.3), Inches(2.2), Inches(3.6), Inches(4), "우리 2,400 vs 스윗아워 13,000\n\n인지도에서 상권 최저권.\n인플루언서 협업으로\n격차 축소가 시급.", 16, True, REDC)

# --- 슬라이드 5: 우리 차별화 ---
s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); title_bar(s, "우리 차별화 포인트")
bullets(s, Inches(0.7), Inches(1.6), Inches(12), Inches(5), [
    ("흑임자 라떼 시그니처 + 바스크 세트로 디저트 강자 C와 정면 대응", "객단가↑ · 시그니처 정체성 강화"),
    ("'조용한 작업 카페' 포지셔닝 전면화", "콘센트·1인 바석·차분함 → A(좌석↓)·C(시끄러움) 약점 공략"),
    ("인스타 인플루언서 협업으로 인지도 격차 축소", "결 맞는 성수 감성/작업카페 인플루언서 3~5명"),
], size=18)

# --- 슬라이드 6: 추천 액션 ---
s = prs.slides.add_slide(BLANK); add_bg(s, WHITE); title_bar(s, "추천 액션")
actions = [
    ("1. 디저트 세트 출시", "흑임자 라떼 + 바스크 세트 → 객단가↑, 디저트 강자 C 견제"),
    ("2. 인플루언서 3~5명 협업", "인지도↑ (마케팅 예산 50만원 내 집행)"),
    ("3. 피크타임 대기·주문 동선 개선", "상권 공통 약점이자 VoC 불만 1위(대기시간) 해소 → 회전율·만족도↑"),
]
top = Inches(1.6)
for i, (t, sub) in enumerate(actions):
    card = s.shapes.add_shape(1, Inches(0.7), top + Inches(i*1.7), Inches(11.9), Inches(1.4))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xFB,0xF3,0xE6); card.line.color.rgb = AMBERC; card.line.width = Pt(1)
    tf = card.text_frame; tf.margin_left = Inches(0.3); tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = t
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = AMBERC; r.font.name = "맑은 고딕"
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
    r2.font.size = Pt(14); r2.font.color.rgb = NAVYC; r2.font.name = "맑은 고딕"

pptx_path = os.path.join(HERE, "하버스카페_경쟁사분석.pptx")
prs.save(pptx_path)
print("PPT 저장:", pptx_path)
print("슬라이드 수:", len(prs.slides._sldIdLst))
