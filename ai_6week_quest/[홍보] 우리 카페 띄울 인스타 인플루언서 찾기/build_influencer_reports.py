# -*- coding: utf-8 -*-
"""하버스카페 인플루언서 섭외안 — influencers.csv → 엑셀(수식·차트·서식) + PPT."""
import csv, os
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
CSV = os.path.join(HERE, "influencers.csv")
XLSX = os.path.join(HERE, "하버스카페_인플루언서_섭외안.xlsx")
PPTX = os.path.join(HERE, "하버스카페_인플루언서_섭외안.pptx")

NAVY = "1F3864"; NAVY2 = "2B4A82"; ORANGE = "C55A11"; LIGHT = "EEF2F8"
RED = "C0392B"; GREEN = "27804E"; GREY = "6E7887"

# 최종 선정(핸들 기준)
PICKS = ["@jakeob.diary", "@minimal.mood.cafe", "@quiet.coffee.kr"]

def load():
    with open(CSV, encoding="utf-8") as f:
        return list(csv.DictReader(f))

rows = load()

# ---------------- EXCEL ----------------
def build_xlsx():
    wb = Workbook()
    thin = Side(style="thin", color="D2D8E2")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    # 시트1: 후보전체
    ws = wb.active; ws.title = "후보전체"
    ws["A1"] = "하버스카페 — 인플루언서 후보 스코어링"
    ws["A1"].font = Font(bold=True, size=15, color=NAVY)
    ws["A2"] = "원천: influencers.csv · 후보 10명 · 성수 카페/작업/디저트 니치"
    ws["A2"].font = Font(size=10, color=GREY)

    headers = ["계정", "분야", "팔로워", "인게이지먼트율(%)", "적합도(/10)", "게시물단가(만원)", "티어",
               "브랜드핏", "타겟일치", "실질도달", "효율(팔/만원)", "종합점수"]
    hr = 4
    for c, h in enumerate(headers, 1):
        cell = ws.cell(hr, c, h)
        cell.font = Font(bold=True, color="FFFFFF", size=10)
        cell.fill = PatternFill("solid", fgColor=NAVY)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border
    ws.row_dimensions[hr].height = 30

    for i, r in enumerate(rows):
        rr = hr + 1 + i
        ws.cell(rr, 1, r["핸들"])                                    # 계정
        ws.cell(rr, 2, r["분야"])
        ws.cell(rr, 3, int(r["팔로워"]))
        ws.cell(rr, 4, float(r["참여율"]))                          # 인게이지먼트율
        ws.cell(rr, 5, f"=H{rr}+I{rr}")                             # 적합도 = 브랜드핏+타겟일치
        ws.cell(rr, 6, int(r["게시물단가"]))
        ws.cell(rr, 7, r["티어"])
        ws.cell(rr, 8, int(r["브랜드핏"]))
        ws.cell(rr, 9, int(r["타겟일치"]))
        # 실질도달 = 팔로워 * 참여율 / 100
        ws.cell(rr, 10, f"=ROUND(C{rr}*D{rr}/100,0)")
        # 효율 = 팔로워 / 단가(F)
        ws.cell(rr, 11, f"=ROUND(C{rr}/F{rr},0)")
        # 종합점수 = 브랜드핏*10 + 타겟일치*10 + 참여율*2
        ws.cell(rr, 12, f"=ROUND(H{rr}*10+I{rr}*10+D{rr}*2,1)")
        is_pick = r["핸들"] in PICKS
        for c in range(1, 13):
            cell = ws.cell(rr, c); cell.border = border
            cell.alignment = Alignment(horizontal="center" if c not in (1, 2) else "left", vertical="center")
            if is_pick:
                cell.fill = PatternFill("solid", fgColor="FCF4EC")
                if c == 1:
                    cell.font = Font(bold=True, color=ORANGE)

    widths = [18, 16, 10, 15, 10, 14, 9, 9, 9, 10, 10, 11]
    for c, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(c)].width = w

    note = hr + 1 + len(rows) + 1
    ws.cell(note, 1, "✅ 주황 하이라이트 = 최종 선정 3명 (합계 50만원). 적합도 = 브랜드핏+타겟일치(/10) · 종합점수 = 브랜드핏×10 + 타겟일치×10 + 참여율×2")
    ws.cell(note, 1).font = Font(size=9, color=ORANGE, italic=True)

    # 차트: 종합점수
    chart = BarChart(); chart.type = "col"; chart.title = "후보별 종합점수"
    chart.height = 8; chart.width = 20
    data = Reference(ws, min_col=12, min_row=hr, max_row=hr + len(rows))
    cats = Reference(ws, min_col=1, min_row=hr + 1, max_row=hr + len(rows))
    chart.add_data(data, titles_from_data=True); chart.set_categories(cats)
    chart.legend = None
    ws.add_chart(chart, f"A{note + 2}")

    # 시트2: 선정·예산
    ws2 = wb.create_sheet("선정·예산")
    ws2["A1"] = "최종 선정 & 예산 시뮬레이션"; ws2["A1"].font = Font(bold=True, size=14, color=NAVY)
    ws2["A2"] = "예산 상한 50만원 · 마이크로/나노 3명 분산 전략"; ws2["A2"].font = Font(size=10, color=GREY)

    h2 = ["핸들", "분야", "팔로워", "참여율(%)", "단가(만원)", "실질도달"]
    hr2 = 4
    for c, h in enumerate(h2, 1):
        cell = ws2.cell(hr2, c, h)
        cell.font = Font(bold=True, color="FFFFFF"); cell.fill = PatternFill("solid", fgColor=NAVY2)
        cell.alignment = Alignment(horizontal="center"); cell.border = border
    pick_rows = [r for r in rows if r["핸들"] in PICKS]
    for i, r in enumerate(pick_rows):
        rr = hr2 + 1 + i
        ws2.cell(rr, 1, r["핸들"]); ws2.cell(rr, 2, r["분야"])
        ws2.cell(rr, 3, int(r["팔로워"])); ws2.cell(rr, 4, float(r["참여율"]))
        ws2.cell(rr, 5, int(r["게시물단가"]))
        ws2.cell(rr, 6, f"=ROUND(C{rr}*D{rr}/100,0)")
        for c in range(1, 7):
            ws2.cell(rr, c).border = border
            ws2.cell(rr, c).alignment = Alignment(horizontal="center" if c != 2 else "left")
    tot = hr2 + 1 + len(pick_rows)
    ws2.cell(tot, 1, "합계"); ws2.cell(tot, 1).font = Font(bold=True)
    ws2.cell(tot, 3, f"=SUM(C{hr2+1}:C{tot-1})")
    ws2.cell(tot, 4, f"=ROUND(AVERAGE(D{hr2+1}:D{tot-1}),1)")
    ws2.cell(tot, 5, f"=SUM(E{hr2+1}:E{tot-1})")
    ws2.cell(tot, 6, f"=SUM(F{hr2+1}:F{tot-1})")
    for c in range(1, 7):
        ws2.cell(tot, c).fill = PatternFill("solid", fgColor=LIGHT)
        ws2.cell(tot, c).font = Font(bold=True); ws2.cell(tot, c).border = border
        ws2.cell(tot, c).alignment = Alignment(horizontal="center" if c != 1 else "left")

    # 예산 시뮬
    b = tot + 2
    ws2.cell(b, 1, "예산 시뮬레이션"); ws2.cell(b, 1).font = Font(bold=True, color=NAVY)
    sim = [("총 예산(만원)", 50), ("집행(선정 3명)", f"=E{tot}"), ("잔액", f"=50-E{tot}")]
    for i, (k, v) in enumerate(sim):
        ws2.cell(b + 1 + i, 1, k); ws2.cell(b + 1 + i, 2, v)
        ws2.cell(b + 1 + i, 1).border = border; ws2.cell(b + 1 + i, 2).border = border

    # 매크로 단독 비교
    cmp = b + 5
    ws2.cell(cmp, 1, "효율 비교: 마이크로·나노 3명 vs 매크로 1명"); ws2.cell(cmp, 1).font = Font(bold=True, color=NAVY)
    ws2.cell(cmp + 1, 1, "전략"); ws2.cell(cmp + 1, 2, "예산(만원)"); ws2.cell(cmp + 1, 3, "실질도달")
    for c in range(1, 4):
        ws2.cell(cmp + 1, c).fill = PatternFill("solid", fgColor=NAVY2)
        ws2.cell(cmp + 1, c).font = Font(bold=True, color="FFFFFF")
        ws2.cell(cmp + 1, c).border = border; ws2.cell(cmp + 1, c).alignment = Alignment(horizontal="center")
    ws2.cell(cmp + 2, 1, "선정 3명(마이크로·나노)"); ws2.cell(cmp + 2, 2, f"=E{tot}"); ws2.cell(cmp + 2, 3, f"=F{tot}")
    ws2.cell(cmp + 3, 1, "@seoul.latteart 단독(매크로)"); ws2.cell(cmp + 3, 2, 60); ws2.cell(cmp + 3, 3, "=ROUND(65000*2.8/100,0)")
    for rr in (cmp + 2, cmp + 3):
        for c in range(1, 4):
            ws2.cell(rr, c).border = border
            ws2.cell(rr, c).alignment = Alignment(horizontal="center" if c != 1 else "left")
    ws2.cell(cmp + 4, 1, "→ 같은 예산으로 선정안이 실질도달 약 1.7배. 참여율·분산 노출 우위.")
    ws2.cell(cmp + 4, 1).font = Font(size=9, color=ORANGE, italic=True)

    for col, w in zip("ABCDEF", [26, 18, 12, 11, 11, 11]):
        ws2.column_dimensions[col].width = w

    wb.save(XLSX)
    print("엑셀 저장:", XLSX)

# ---------------- PPT ----------------
def _rgb(h): return RGBColor.from_string(h)
def _fill(shape, hexc): shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(hexc); shape.line.fill.background()
def _tb(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = _rgb(color); r.font.name = "맑은 고딕"
    return tb

def build_pptx():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W = 13.333

    # 1 표지
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height); _fill(bg, NAVY)
    bar = s.shapes.add_shape(1, 0, Inches(3.0), prs.slide_width, Inches(0.09)); _fill(bar, ORANGE)
    _tb(s, 1.0, 3.2, 11, 1.2, "하버스카페 인플루언서 섭외안", 40, "FFFFFF", bold=True)
    _tb(s, 1.05, 4.4, 11, 0.8, "성수동 · 인스타 인지도 확대 → 예산 50만원 실행안", 18, "C8D2E4")
    _tb(s, 1.05, 6.4, 11, 0.5, "원천: influencers.csv · 경쟁사 분석 추천 액션 2번", 12, "9AA6BE")

    # 2 타겟 정의
    s = prs.slides.add_slide(blank)
    _tb(s, 0.6, 0.4, 12, 0.9, "타겟 정의", 30, NAVY, bold=True)
    tgt = [("우리 카페 핵심 손님", "성수동·서울숲 25~35세 직장인·프리랜서·디자이너 · 평일 노트북 작업(콘센트·조용함) · 흑임자 라떼·디저트에 반응 · 인스타로 카페 검색·저장"),
           ("캠페인 도달 타겟", "위 손님과 팔로워가 겹치는 '작업카페·감성/미니멀 카페·조용한 카페' 콘텐츠 소비자 → 이 타겟과 겹치는 인플루언서 우선"),
           ("왜 이 타겟인가", "우리 강점(조용한 작업환경 + 흑임자 시그니처)이 정확히 이 층에 소구 → 방문·저장 전환율↑, 경쟁사 약점(좌석↓·시끄러움) 공략점과 일치")]
    y = 1.55
    for t, d in tgt:
        card = s.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(12.1), Inches(1.5)); _fill(card, LIGHT)
        tab = s.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(0.09), Inches(1.5)); _fill(tab, ORANGE)
        _tb(s, 0.95, y + 0.16, 11, 0.6, t, 20, ORANGE, bold=True)
        _tb(s, 0.95, y + 0.62, 11.4, 0.8, d, 14, "212836")
        y += 1.7

    # 3 선정 기준
    s = prs.slides.add_slide(blank)
    _tb(s, 0.6, 0.4, 12, 0.9, "선정 기준 — 팔로워 수 ≠ 효과", 30, NAVY, bold=True)
    crit = [("① 브랜드 핏", "조용한 작업카페·흑임자·디저트 감성과 결이 맞는가"),
            ("② 참여율", "좋아요·댓글·저장률 — 나노·마이크로가 매크로보다 3~5배 높음"),
            ("③ 타겟 일치", "성수 직장인·프리랜서·작업 손님과 팔로워가 겹치는가"),
            ("④ 예산 효율", "게시물 단가 대비 실질 도달(팔로워×참여율)")]
    y = 1.5
    for t, d in crit:
        card = s.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(12.1), Inches(1.1)); _fill(card, LIGHT)
        _tb(s, 0.9, y + 0.15, 4, 0.7, t, 22, ORANGE, bold=True)
        _tb(s, 4.6, y + 0.2, 8.0, 0.7, d, 16, "212836")
        y += 1.28

    # 3 후보 롱리스트
    s = prs.slides.add_slide(blank)
    _tb(s, 0.6, 0.4, 12, 0.9, "후보 롱리스트 (10명) — 계정·팔로워·인게이지먼트율·적합도", 24, NAVY, bold=True)
    cols = ["계정", "분야", "팔로워", "인게이지먼트율", "적합도", "단가"]
    top = rows
    tbl = s.shapes.add_table(len(top) + 1, len(cols), Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4)).table
    cw = [2.5, 2.9, 1.8, 2.3, 1.4, 1.2]
    for i, w in enumerate(cw): tbl.columns[i].width = Inches(w)
    for c, h in enumerate(cols):
        cell = tbl.cell(0, c); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(NAVY2)
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
        pr.runs[0].font.size = Pt(13); pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = _rgb("FFFFFF")
    for r, row in enumerate(top, 1):
        pick = row["핸들"] in PICKS
        fit = int(row["브랜드핏"]) + int(row["타겟일치"])
        vals = [row["핸들"], row["분야"], f'{int(row["팔로워"]):,}', f'{row["참여율"]}%',
                f'{fit}/10', f'{row["게시물단가"]}만']
        for c, v in enumerate(vals):
            cell = tbl.cell(r, c); cell.text = str(v)
            cell.fill.solid(); cell.fill.fore_color.rgb = _rgb("FCF4EC" if pick else ("F5F7FB" if r % 2 else "FFFFFF"))
            pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER if c not in (0, 1) else PP_ALIGN.LEFT
            pr.runs[0].font.size = Pt(11)
            pr.runs[0].font.color.rgb = _rgb(ORANGE if (pick and c == 0) else "212836")
            pr.runs[0].font.bold = pick and c == 0

    # 4 최종 선정
    s = prs.slides.add_slide(blank)
    _tb(s, 0.6, 0.4, 12, 0.9, "✅ 최종 선정 3명 — 합계 50만원", 30, NAVY, bold=True)
    picks = [r for r in rows if r["핸들"] in PICKS]
    cols4 = ["핸들", "분야", "팔로워", "참여율", "단가", "선정 이유"]
    reasons = {"@jakeob.diary": "'작업카페' 포지션 정확 일치·높은 참여율",
               "@minimal.mood.cafe": "조용·차분 무드 = 우리 인테리어 결",
               "@quiet.coffee.kr": "나노 최고 참여율·흑임자 저장 유도"}
    tbl = s.shapes.add_table(len(picks) + 2, len(cols4), Inches(0.6), Inches(1.5), Inches(12.1), Inches(3.2)).table
    for i, w in enumerate([2.3, 2.3, 1.5, 1.3, 1.2, 3.5]): tbl.columns[i].width = Inches(w)
    for c, h in enumerate(cols4):
        cell = tbl.cell(0, c); cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(NAVY2)
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
        pr.runs[0].font.size = Pt(13); pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = _rgb("FFFFFF")
    for r, row in enumerate(picks, 1):
        vals = [row["핸들"], row["분야"], f'{int(row["팔로워"]):,}', f'{row["참여율"]}%',
                f'{row["게시물단가"]}만', reasons[row["핸들"]]]
        for c, v in enumerate(vals):
            cell = tbl.cell(r, c); cell.text = str(v); cell.fill.solid(); cell.fill.fore_color.rgb = _rgb("FCF4EC")
            pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER if c not in (0, 1, 5) else PP_ALIGN.LEFT
            pr.runs[0].font.size = Pt(11.5); pr.runs[0].font.color.rgb = _rgb(ORANGE if c == 0 else "212836"); pr.runs[0].font.bold = c == 0
    # 합계행
    rr = len(picks) + 1
    tot_vals = ["합계", "", f'{sum(int(p["팔로워"]) for p in picks):,}', "평균 7.6%", "50만", "실질도달 ≈ 3,150"]
    for c, v in enumerate(tot_vals):
        cell = tbl.cell(rr, c); cell.text = str(v) if v else " "; cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(NAVY)
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER if c not in (0, 5) else PP_ALIGN.LEFT
        if pr.runs:
            pr.runs[0].font.size = Pt(12); pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = _rgb("FFFFFF")
    _tb(s, 0.6, 5.1, 12, 1.6,
        "왜 이 조합? @seoul.latteart 65,000명 단독(60만·참여율 2.8%)은 예산 초과 + 실질 인게이지먼트 낮음.\n"
        "선정 3명 실질도달 ≈ 3,150 vs 매크로 단독 ≈ 1,820 → 같은 예산으로 1.7배 효율 + 3개 니치 분산 노출.",
        15, "212836")

    # 5 Top3 컨택 메시지 초안
    s = prs.slides.add_slide(blank)
    _tb(s, 0.6, 0.4, 12, 0.9, "Top 3 컨택 메시지 초안 (맞춤 DM)", 28, NAVY, bold=True)
    dms = [
        ("@jakeob.diary",
         "안녕하세요 :) 성수동 '하버스카페'입니다. '작업일기' 피드 잘 보고 있어요 — 콘센트·1인 바석 갖춘 조용한 작업카페라 결이 잘 맞을 것 같아 연락드려요. 흑임자 라떼+바스크 세트 대접하고, 작업 브이로그(피드1+릴스1) 협업 제안드립니다. 소정의 원고료도 준비돼 있어요!"),
        ("@minimal.mood.cafe",
         "안녕하세요 :) 미니멀하고 차분한 무드 피드에 반해 연락드려요. 저희 인테리어도 조용·미니멀 결이라 잘 어울릴 것 같습니다. 흑임자 라떼+바스크 세트와 함께 감성 컷 콘텐츠 협업 제안드려요. 편하실 때 회신 부탁드립니다!"),
        ("@quiet.coffee.kr",
         "안녕하세요 :) '조용한 카페' 큐레이션 늘 잘 보고 있어요. 노트북 작업 집중 잘 되는 저희 하버스카페도 딱 맞을 것 같아 연락드립니다. 흑임자 라떼+바스크 세트 대접하고, 저장 유도되는 무드 릴스 협업 제안드려요!"),
    ]
    y = 1.5
    for handle, msg in dms:
        card = s.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(12.1), Inches(1.75)); _fill(card, "F5F7FB")
        tab = s.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(0.09), Inches(1.75)); _fill(tab, ORANGE)
        _tb(s, 0.9, y + 0.12, 11, 0.5, handle, 17, ORANGE, bold=True)
        _tb(s, 0.9, y + 0.55, 11.5, 1.1, msg, 12.5, "212836")
        y += 1.92

    # 6 다음 액션
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height); _fill(bg, NAVY)
    bar = s.shapes.add_shape(1, 0, Inches(1.5), prs.slide_width, Inches(0.06)); _fill(bar, ORANGE)
    _tb(s, 0.7, 0.55, 12, 0.9, "다음 액션", 34, "FFFFFF", bold=True)
    acts = [("1. DM 발송", "3명에게 맞춤 DM(바터+소액 원고료) — 흑임자 라떼+바스크 세트 제공"),
            ("2. 콘텐츠 협의", "피드 1 + 릴스/스토리 1 · '조용한 작업카페·흑임자 라떼' 키워드·위치 태그"),
            ("3. 성과 측정", "게시 2주 후 저장 수·프로필 방문·'인스타 보고 왔어요' 방문자 집계")]
    y = 1.9
    for t, d in acts:
        card = s.shapes.add_shape(1, Inches(0.7), Inches(y), Inches(11.9), Inches(1.35)); _fill(card, "F5F7FB")
        _tb(s, 1.0, y + 0.18, 4, 0.7, t, 22, NAVY, bold=True)
        _tb(s, 4.7, y + 0.24, 7.6, 0.9, d, 15, "46505F")
        y += 1.55

    prs.save(PPTX)
    print("PPT 저장:", PPTX, "· 슬라이드 수:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    print("=== 하버스카페 인플루언서 섭외안 → 엑셀·PPT 생성 ===")
    build_xlsx()
    build_pptx()
    print("완료")
